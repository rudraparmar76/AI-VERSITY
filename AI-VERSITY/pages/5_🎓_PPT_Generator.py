import streamlit as st 
import google.generativeai as genai
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re
import requests
from pptx.util import Inches
from pptx import Presentation
from api import Gemini_Key
import os
import pandas as pd
from datetime import datetime

if 'signedin' in st.session_state and st.session_state.signedin:
    st.set_page_config(page_title="SlideCraft", page_icon="🎓",layout="wide")
    st.markdown("# SlideCraft 🎓")
    st.sidebar.header(" SlideCraft 🎓")

    usage_data_path = r'UsageData/pptgen_usage_data.csv'
    if os.path.exists(usage_data_path):
        usage_data = pd.read_csv(usage_data_path)
    else:
        usage_data = pd.DataFrame(columns=['Prompt','slides','Timestamp'])


    genai.configure(api_key=Gemini_Key)

    # def fetch_image_from_unplash(topic):
    #     acess_key = "Iwk_YMtik0GEgibRGYx05pWIZVIvQzdegmZ9oKXwK3U"
    #     secret_key = "Z88t8Af5Ej8Iw0EKeEYkyo1FNRk57uwPqLkxcn9sIww"
    #     headers = {
    #         "Authorization":f"Client-ID {acess_key}"
    #     }
    #     url = f"https://api.unsplash.com/search/photos?query={topic}&per_page=1"
    #     response = requests.get(url, headers=headers)
    #     if response.status_code == 200:
    #         data = response.json()
    #         if data["total"] > 0:
    #             return data["results"][0]["urls"]["small"]
    #     return None


    def get_gemini_response(prompt):
        model = genai.GenerativeModel('gemini-pro')
        response = model.generate_content(prompt)
        return response

    sub_titles = []
    def refine_subtopics(sub_topics, sub_titles):
        for sub_topic in sub_topics:
            sub_titles.append(sub_topic[3:].replace('"',""))
        return sub_titles

    content = []
    def content_generation(sub_titles):
        for i in sub_titles:
            prompt = f"Generate a content of {i} for presentation slide on the 2 bullet point only each of point 20 tokens"
            model = genai.GenerativeModel('gemini-pro')
            response = model.generate_content(prompt)
            content.append(response.text)
        return content

    def clean_text(text):
        # Removing extra whitespaces and newlines
        cleaned_text = re.sub('\s+', ' ', text).strip()

        # Removing markdown-style bullet points, asterisks, and numeric bullet points
        cleaned_text = re.sub(r'[*-]\s*|\d+\.\s*', '', cleaned_text)

        # Removing extra spaces before and after colons
        cleaned_text = re.sub(r'\s*:\s*', ': ', cleaned_text)

        # Removing extra spaces before and after hyphens
        cleaned_text = re.sub(r'\s*-\s*', ' - ', cleaned_text)

        return cleaned_text


    def split_sentences(text):
        # Split the text into sentences using regular expression
        sentences = re.split(r'(?<=\.)\s+', text)

        # Capitalizing the first letter of each sentence
        sentences = [sentence.capitalize() for sentence in sentences]

        return sentences
    def replace_and_capitalize(text):
        # Defining a function to replace and capitalize the text between colons
        def replace_and_capitalize_colon(match):
            return match.group(1) + match.group(2).capitalize() + match.group(3)

        # Using regular expression to find and replace text between colons
        result = re.sub(r'(:\s*)(.*?)(\s*:[^:]|$)', replace_and_capitalize_colon, text)

        return result


    final_content = []
    def refine_final_content(content):
        for i in content:
            cleaned_text = clean_text(i)
            sentences = split_sentences(cleaned_text)
            final_content.append(sentences)
        print("final content ready....")
        return final_content

    powerpoint = Presentation()

    def slide_maker(powerpoint, topic,sub_titles, final_content):
        title_slide_layout = powerpoint.slide_layouts[0]
        title_slide = powerpoint.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        title.text = topic
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        content = title_slide.placeholders[1]
        content.text = "Created By AI Gemini Model"
        for i in range(len(sub_titles)):
            bulletLayout = powerpoint.slide_layouts[1]
            secondSlide = powerpoint.slides.add_slide(bulletLayout)
            # accessing the attributes of shapes
            myShapes = secondSlide.shapes
            titleShape = myShapes.title
            bodyShape = myShapes.placeholders[1]
            titleShape.text = sub_titles[i]
            titleShape.text_frame.paragraphs[0].font.size = Pt(24)
            titleShape.text_frame.paragraphs[0].font.bold = True
            tFrame = bodyShape.text_frame
            print("Topic Generated")
            for point in final_content[i]:
                point = re.sub(r':[^:]+:', ':', point)
                point = replace_and_capitalize(point)
                p = tFrame.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
                p.space_after = Pt(10)
        return powerpoint
    def download_button(file_path,topic):
        # Reading the content of the PPT file
        with open(file_path, "rb") as file:
            ppt_content = file.read()

        # Creating a download button for the PPT file
        st.download_button(
            label="Download PowerPoint",
            data=ppt_content,
            file_name=f"{topic}.pptx", 
            key="ppt_download_button"
        )


    topic=st.text_area("Input Prompt: ",key="input",placeholder="Enter the topic for which you want to generate the presentation")
    no_of_slide=st.text_input("Enter Number Of Slide: ",key="slide")

    submit=st.button("Generate")
    if submit:
        with st.spinner('Generating Presentation...'):
            prompt =f"Generate a {no_of_slide} sub-titles only  on the topic of {topic}"
            response = get_gemini_response(prompt)
            print("Topic Generated")
            sub_topics = response.text.split("\n")
            sub_titles = refine_subtopics(sub_topics, sub_titles)
            print("Sub Titles")
            content = content_generation(sub_titles)
            print("content Generated")
            final_content = refine_final_content(content)
            #cleaned_text = clean_text(content[0])
            #sentences = split_sentences(cleaned_text)
            print("final content ready")
            powerpoint = slide_maker(powerpoint,topic, sub_titles, final_content)
            powerpoint.save(f"Powerpoint/{topic}.pptx")
            st.text("Presentation Ready")
            download_button(f"Powerpoint/{topic}.pptx",topic)
            print("Presentation Ready")
            usage_data = usage_data._append({'Prompt': topic,'slides':no_of_slide,'Timestamp': datetime.now()}, ignore_index=True)
            usage_data.to_csv(usage_data_path, index=False)
else:
    st.warning("Please sign in to generate the presentation.")